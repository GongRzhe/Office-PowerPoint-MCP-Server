"""
Presentation management utilities for PowerPoint MCP Server.
Functions for creating, opening, saving, and managing presentations.
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from typing import Dict, List, Optional
import os


# Standard slide dimensions in EMUs (English Metric Units)
# 914400 EMUs = 1 inch
SLIDE_DIMENSIONS = {
    "4:3": {
        "width": Inches(10),      # 10 inches
        "height": Inches(7.5),    # 7.5 inches
        "name": "Standard (4:3)"
    },
    "16:9": {
        "width": Inches(13.333),  # 13.333 inches
        "height": Inches(7.5),    # 7.5 inches
        "name": "Widescreen (16:9)"
    },
    "16:10": {
        "width": Inches(10),      # 10 inches
        "height": Inches(6.25),   # 6.25 inches
        "name": "Widescreen (16:10)"
    },
    "a4": {
        "width": Inches(10.833),  # A4 portrait width
        "height": Inches(7.5),    # A4 portrait height
        "name": "A4 Paper"
    }
}


def create_presentation(aspect_ratio: str = "16:9") -> Presentation:
    """
    Create a new PowerPoint presentation with specified aspect ratio.

    Args:
        aspect_ratio: Slide aspect ratio. Options: "4:3", "16:9", "16:10", "a4".
                     Default is "16:9" for widescreen.

    Returns:
        A new Presentation object with the specified slide dimensions
    """
    pres = Presentation()

    # Normalize the aspect ratio string
    aspect_ratio = aspect_ratio.lower().strip()

    # Get dimensions for the specified aspect ratio
    if aspect_ratio in SLIDE_DIMENSIONS:
        dimensions = SLIDE_DIMENSIONS[aspect_ratio]
        pres.slide_width = dimensions["width"]
        pres.slide_height = dimensions["height"]
    else:
        # Default to 16:9 if invalid ratio provided
        dimensions = SLIDE_DIMENSIONS["16:9"]
        pres.slide_width = dimensions["width"]
        pres.slide_height = dimensions["height"]

    return pres


def open_presentation(file_path: str) -> Presentation:
    """
    Open an existing PowerPoint presentation.
    
    Args:
        file_path: Path to the PowerPoint file
        
    Returns:
        A Presentation object
    """
    return Presentation(file_path)


def create_presentation_from_template(template_path: str) -> Presentation:
    """
    Create a new PowerPoint presentation from a template file.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        A new Presentation object based on the template
        
    Raises:
        FileNotFoundError: If the template file doesn't exist
        Exception: If the template file is corrupted or invalid
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    if not template_path.lower().endswith(('.pptx', '.potx')):
        raise ValueError("Template file must be a .pptx or .potx file")
    
    try:
        # Load the template file as a presentation
        presentation = Presentation(template_path)
        return presentation
    except Exception as e:
        raise Exception(f"Failed to load template file '{template_path}': {str(e)}")


def save_presentation(presentation: Presentation, file_path: str) -> str:
    """
    Save a PowerPoint presentation to a file.
    
    Args:
        presentation: The Presentation object
        file_path: Path where the file should be saved
        
    Returns:
        The file path where the presentation was saved
    """
    presentation.save(file_path)
    return file_path


def get_template_info(template_path: str) -> Dict:
    """
    Get information about a template file.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        Dictionary containing template information
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    try:
        presentation = Presentation(template_path)
        
        # Get slide layouts
        layouts = get_slide_layouts(presentation)
        
        # Get core properties
        core_props = get_core_properties(presentation)
        
        # Get slide count
        slide_count = len(presentation.slides)
        
        # Get file size
        file_size = os.path.getsize(template_path)
        
        return {
            "template_path": template_path,
            "file_size_bytes": file_size,
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props
        }
    except Exception as e:
        raise Exception(f"Failed to read template info from '{template_path}': {str(e)}")


def get_presentation_info(presentation: Presentation) -> Dict:
    """
    Get information about a presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Dictionary containing presentation information
    """
    try:
        # Get slide layouts
        layouts = get_slide_layouts(presentation)
        
        # Get core properties
        core_props = get_core_properties(presentation)
        
        # Get slide count
        slide_count = len(presentation.slides)
        
        return {
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props,
            "slide_width": presentation.slide_width,
            "slide_height": presentation.slide_height
        }
    except Exception as e:
        raise Exception(f"Failed to get presentation info: {str(e)}")


def get_slide_layouts(presentation: Presentation) -> List[Dict]:
    """
    Get all available slide layouts in the presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        A list of dictionaries with layout information
    """
    layouts = []
    for i, layout in enumerate(presentation.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders)
        }
        layouts.append(layout_info)
    return layouts


def set_core_properties(presentation: Presentation, title: str = None, subject: str = None,
                       author: str = None, keywords: str = None, comments: str = None) -> None:
    """
    Set core document properties.
    
    Args:
        presentation: The Presentation object
        title: Document title
        subject: Document subject
        author: Document author
        keywords: Document keywords
        comments: Document comments
    """
    core_props = presentation.core_properties
    
    if title is not None:
        core_props.title = title
    if subject is not None:
        core_props.subject = subject
    if author is not None:
        core_props.author = author
    if keywords is not None:
        core_props.keywords = keywords
    if comments is not None:
        core_props.comments = comments


def get_core_properties(presentation: Presentation) -> Dict:
    """
    Get core document properties.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Dictionary containing core properties
    """
    core_props = presentation.core_properties
    
    return {
        "title": core_props.title,
        "subject": core_props.subject,
        "author": core_props.author,
        "keywords": core_props.keywords,
        "comments": core_props.comments,
        "created": core_props.created.isoformat() if core_props.created else None,
        "last_modified_by": core_props.last_modified_by,
        "modified": core_props.modified.isoformat() if core_props.modified else None
    }