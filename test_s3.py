#!/usr/bin/env python
"""
Test script for S3 storage functionality in PowerPoint MCP Server.
Tests all S3 tools with a local MinIO server.
"""
import sys
import os

# Add the current directory to the path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from utils.s3_utils import (
    create_s3_client,
    upload_presentation_to_s3,
    download_presentation_from_s3,
    list_s3_objects,
    delete_s3_object,
    get_s3_object_metadata
)
from utils.presentation_utils import create_presentation
from pptx.util import Inches
import boto3
from botocore.exceptions import ClientError


def print_section(title):
    """Print a section header."""
    print("\n" + "=" * 60)
    print(f"  {title}")
    print("=" * 60)


def test_s3_functionality():
    """Test all S3 functionality."""
    
    # Configuration
    ENDPOINT_URL = "http://localhost:9000"
    ACCESS_KEY = "minioadmin"
    SECRET_KEY = "minioadmin"
    BUCKET_NAME = "test-bucket-powerpoint"
    OBJECT_KEY = "test-presentations/sample.pptx"
    
    print_section("S3 Storage Test Suite")
    print(f"Endpoint: {ENDPOINT_URL}")
    print(f"Bucket: {BUCKET_NAME}")
    print(f"Object Key: {OBJECT_KEY}")
    
    # Step 1: Create S3 client
    print_section("Step 1: Creating S3 Client")
    try:
        s3_client = create_s3_client(
            endpoint_url=ENDPOINT_URL,
            access_key=ACCESS_KEY,
            secret_key=SECRET_KEY,
            region="us-east-1"
        )
        print("✓ S3 client created successfully")
    except Exception as e:
        print(f"✗ Failed to create S3 client: {e}")
        return False
    
    # Step 2: Create bucket
    print_section("Step 2: Creating Test Bucket")
    try:
        s3_client.create_bucket(Bucket=BUCKET_NAME)
        print(f"✓ Bucket '{BUCKET_NAME}' created successfully")
    except ClientError as e:
        if e.response['Error']['Code'] == 'BucketAlreadyOwnedByYou':
            print(f"✓ Bucket '{BUCKET_NAME}' already exists")
        else:
            print(f"✗ Failed to create bucket: {e}")
            return False
    
    # Step 3: Create a test presentation
    print_section("Step 3: Creating Test Presentation")
    try:
        pres = create_presentation()
        
        # Add a slide with some content
        slide_layout = pres.slide_layouts[0]  # Title slide
        slide = pres.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = "S3 Storage Test"
        
        # Add subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Testing PowerPoint MCP Server S3 Integration"
        
        # Add another slide with content
        slide_layout = pres.slide_layouts[1]  # Title and content
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Test Content"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "This presentation was created and uploaded to S3 storage."
        
        print(f"✓ Created presentation with {len(pres.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to create presentation: {e}")
        return False
    
    # Step 4: Upload presentation to S3
    print_section("Step 4: Uploading Presentation to S3")
    try:
        result = upload_presentation_to_s3(
            presentation=pres,
            s3_client=s3_client,
            bucket_name=BUCKET_NAME,
            object_key=OBJECT_KEY,
            metadata={
                "test": "true",
                "author": "test-script",
                "version": "1.0"
            }
        )
        print(f"✓ Upload successful")
        print(f"  - Bucket: {result['bucket']}")
        print(f"  - Key: {result['key']}")
        print(f"  - Size: {result['size_bytes']} bytes")
    except Exception as e:
        print(f"✗ Failed to upload presentation: {e}")
        return False
    
    # Step 5: List objects in bucket
    print_section("Step 5: Listing Objects in Bucket")
    try:
        result = list_s3_objects(
            s3_client=s3_client,
            bucket_name=BUCKET_NAME,
            prefix="test-presentations/"
        )
        print(f"✓ Found {result['object_count']} object(s)")
        for obj in result['objects']:
            print(f"  - {obj['key']} ({obj['size_bytes']} bytes)")
    except Exception as e:
        print(f"✗ Failed to list objects: {e}")
        return False
    
    # Step 6: Get object metadata
    print_section("Step 6: Getting Object Metadata")
    try:
        result = get_s3_object_metadata(
            s3_client=s3_client,
            bucket_name=BUCKET_NAME,
            object_key=OBJECT_KEY
        )
        print(f"✓ Retrieved metadata")
        print(f"  - Size: {result['size_bytes']} bytes")
        print(f"  - Content Type: {result['content_type']}")
        print(f"  - Last Modified: {result['last_modified']}")
        print(f"  - Custom Metadata: {result['metadata']}")
    except Exception as e:
        print(f"✗ Failed to get metadata: {e}")
        return False
    
    # Step 7: Download presentation from S3
    print_section("Step 7: Downloading Presentation from S3")
    try:
        downloaded_pres = download_presentation_from_s3(
            s3_client=s3_client,
            bucket_name=BUCKET_NAME,
            object_key=OBJECT_KEY
        )
        print(f"✓ Download successful")
        print(f"  - Slides: {len(downloaded_pres.slides)}")
        print(f"  - Slide width: {downloaded_pres.slide_width}")
        print(f"  - Slide height: {downloaded_pres.slide_height}")
        
        # Verify content
        first_slide = downloaded_pres.slides[0]
        title_text = first_slide.shapes.title.text
        print(f"  - First slide title: '{title_text}'")
        
        if title_text == "S3 Storage Test":
            print("✓ Content verification passed")
        else:
            print("✗ Content verification failed")
            return False
            
    except Exception as e:
        print(f"✗ Failed to download presentation: {e}")
        return False
    
    # Step 8: Pause for manual verification
    print_section("Step 8: Manual Verification")
    print(f"\n✓ Presentation uploaded successfully!")
    print(f"\nYou can now verify the file in MinIO:")
    print(f"  1. Open http://localhost:9001 in your browser")
    print(f"  2. Login with: minioadmin / minioadmin")
    print(f"  3. Navigate to bucket: {BUCKET_NAME}")
    print(f"  4. Find object: {OBJECT_KEY}")
    print(f"  5. Download and open the .pptx file to verify content")
    print(f"\nPress Enter to continue with cleanup...")
    input()
    
    # Step 9: Delete object from S3
    print_section("Step 9: Deleting Object from S3")
    try:
        result = delete_s3_object(
            s3_client=s3_client,
            bucket_name=BUCKET_NAME,
            object_key=OBJECT_KEY
        )
        print(f"✓ Delete successful")
        print(f"  - {result['message']}")
    except Exception as e:
        print(f"✗ Failed to delete object: {e}")
        return False
    
    # Step 10: Verify deletion
    print_section("Step 10: Verifying Deletion")
    try:
        result = list_s3_objects(
            s3_client=s3_client,
            bucket_name=BUCKET_NAME,
            prefix="test-presentations/"
        )
        if result['object_count'] == 0:
            print("✓ Object successfully deleted")
        else:
            print(f"✗ Object still exists ({result['object_count']} objects found)")
            return False
    except Exception as e:
        print(f"✗ Failed to verify deletion: {e}")
        return False
    
    # Step 11: Clean up bucket
    print_section("Step 11: Cleaning Up Test Bucket")
    try:
        s3_client.delete_bucket(Bucket=BUCKET_NAME)
        print(f"✓ Bucket '{BUCKET_NAME}' deleted successfully")
    except Exception as e:
        print(f"⚠ Warning: Failed to delete bucket (may not be empty): {e}")
    
    # Success!
    print_section("Test Results")
    print("✓ All tests passed successfully!")
    print("\nS3 storage functionality is working correctly.")
    return True


if __name__ == "__main__":
    print("\n" + "=" * 60)
    print("  PowerPoint MCP Server - S3 Storage Test")
    print("=" * 60)
    print("\nMake sure MinIO is running on http://localhost:9000")
    print("Default credentials: minioadmin / minioadmin\n")
    
    try:
        success = test_s3_functionality()
        
        if success:
            print("\n" + "=" * 60)
            print("  ✓ ALL TESTS PASSED")
            print("=" * 60)
            sys.exit(0)
        else:
            print("\n" + "=" * 60)
            print("  ✗ TESTS FAILED")
            print("=" * 60)
            sys.exit(1)
            
    except KeyboardInterrupt:
        print("\n\nTest interrupted by user")
        sys.exit(1)
    except Exception as e:
        print(f"\n\nUnexpected error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
