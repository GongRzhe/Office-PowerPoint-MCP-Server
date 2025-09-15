# tools/s3_xlsx_tools.py
from typing import Dict, List, Optional
from io import BytesIO

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

def register_s3_xlsx_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb,
):
    """
    既存の register_*_tools と同じシグネチャで登録する。
    """

    # 依存: boto3 / openpyxl を使って S3 から xlsx を読み、現在のプレゼンに
    # 1) 表スライド（省略可） 2) 棒グラフスライド を追加する最小ツール。
    import boto3
    from openpyxl import load_workbook
    from pptx.util import Inches

    s3 = boto3.client("s3")

    @app.tool()
    def build_slides_from_s3_xlsx(
        bucket: str,
        key: str,
        sheet: Optional[str] = None,
        x_col: Optional[str] = None,
        y_cols: Optional[List[str]] = None,
        add_table: bool = True,
        title: str = "From S3 XLSX",
    ) -> Dict:
        """Read an XLSX from S3 and add a table+chart slide to the current presentation."""
        ...

        # ---- 1) S3からxlsxをバイトで取得 ----
        obj = s3.get_object(Bucket=bucket, Key=key)  # 返り値のBodyはStreamingBody
        xlsx_bytes = obj["Body"].read()             # bytesに読み出し（小〜中サイズ前提） :contentReference[oaicite:0]{index=0}

        # ---- 2) メモリからopenpyxlで読み込む ----
        wb = load_workbook(filename=BytesIO(xlsx_bytes), data_only=True)  # BytesIOでOK :contentReference[oaicite:1]{index=1}
        ws = wb[sheet] if sheet else wb.active

        # 1行目をヘッダとしてdict化
        rows = list(ws.values)
        header = [str(h) for h in rows[0]]
        body = rows[1:]
        records = [{h: r[i] for i, h in enumerate(header)} for r in body]

        # ---- 3) 列推測 ----
        x_col = x_col or header[0]
        if not y_cols:
            # 数値っぽい列を自動選定
            def _numlike(v):
                try:
                    float(v); return True
                except: return False
            y_cols = [c for c in header if c != x_col and any(_numlike(rec.get(c)) for rec in records[:8])] or header[1:2]

        # ---- 4) 現在のプレゼン取得 ----
        pres_id = get_current_presentation_id()
        if pres_id not in presentations:
            return {"error": "No presentation loaded. Create or open one first."}
        prs = presentations[pres_id]

        # ---- 5) 表スライド（任意）----
        if add_table:
            layout = prs.slide_layouts[5]  # Title Only 等
            slide_tbl = prs.slides.add_slide(layout)
            slide_tbl.shapes.title.text = f"{title}（表）"

            n_rows = min(len(records) + 1, 20)
            n_cols = min(len(header), 10)
            shape = slide_tbl.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            table = shape.table  # add_tableはshapeを返す → tableはshape.tableで取得 :contentReference[oaicite:2]{index=2}

            for j, col in enumerate(header[:n_cols]):
                table.cell(0, j).text = str(col)
            for i in range(n_rows - 1):
                for j, col in enumerate(header[:n_cols]):
                    table.cell(i + 1, j).text = "" if i >= len(records) else str(records[i].get(col, ""))

        # ---- 6) グラフスライド（棒）----
        layout = prs.slide_layouts[5]
        slide_ch = prs.slides.add_slide(layout)
        slide_ch.shapes.title.text = f"{title}（グラフ）"

        chart_data = CategoryChartData()                   # 生成にもreplaceにも使うChartData系 :contentReference[oaicite:3]{index=3}
        chart_data.categories = [rec.get(x_col, "") for rec in records]
        for yc in y_cols:
            vals = []
            for rec in records:
                v = rec.get(yc)
                try:
                    vals.append(float(v) if v is not None else 0.0)
                except:
                    vals.append(0.0)
            chart_data.add_series(str(yc), vals)

        x, y, cx, cy = Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5)
        slide_ch.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)  # add_chartにChartDataを渡すのが正道 :contentReference[oaicite:4]{index=4}

        return {
            "presentation_id": pres_id,
            "rows": len(records),
            "x_col": x_col,
            "y_cols": y_cols,
        }
